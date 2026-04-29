"""Creator Experts — internal proposal deck for Aaron Langerman (Mercor) — v3.

Framed as an internal proposal from two current Mercor contractors (Logan +
Aarnav, both on the CUA-envs contract for Anthropic under Drew Geoly).

Mercor design tokens mirror work.mercor.com 1:1:
  - bg #ffffff, bg-elev #fafafa
  - fg #0a0a0a, fg-muted #6b7280, fg-subtle #9ca3af
  - border #e5e7eb
  - accent #7857ff (Mercor purple), accent-soft #ede9fe
  - typography: Inter

v3 structure (8 slides) — addresses Aaron's verbatim feedback:
  1. Cover                          — "A high-friction labor market. Where AI compounds."
  2. THE BROKEN WORKFLOW            — 5-step horizontal flow + bottleneck stats
  3. WHY NOW                        — three numbered drivers
  4. SIGNAL                         — comparison ledger + THE MOAT row + 68% Lumanu
  5. VERIFICATION                   — 4-step ladder, steps 3+4 honest "manual review"
  6. DAY-1 REVENUE                  — $170 placement + $0-$50 perf + RL licensing
  7. PROVE BY AUGUST                — 4 falsifiable proof points + verdict gate
  8. NOW → LIVE DEMO                — vercel URL handoff

Run: python3 scripts/build_aaron_deck.py
Output: docs/pitch/Mercor-Creators-Domain.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ─── Tokens ─────────────────────────────────────────────────────────────────
BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_ELEV = RGBColor(0xFA, 0xFA, 0xFA)
FG = RGBColor(0x0A, 0x0A, 0x0A)
FG_MUTED = RGBColor(0x6B, 0x72, 0x80)
FG_SUBTLE = RGBColor(0x9C, 0xA3, 0xAF)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
ACCENT = RGBColor(0x78, 0x57, 0xFF)
ACCENT_SOFT = RGBColor(0xED, 0xE9, 0xFE)

FONT_DISPLAY = "Inter"
FONT_BODY = "Inter"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

REPO_ROOT = Path(__file__).resolve().parents[1]
MERCOR_LOGO = REPO_ROOT / "docs" / "pitch" / "mercor-logo.png"

# Verbatim from frontend/src/lib/data/source-of-truth.ts
DAY1_REVENUE_HEADLINE = (
    "Logan × Celsius signed: $170 placement + $0–$50 perf kicker "
    "+ per-task RL revenue when this becomes a Mercor world."
)


# ─── Primitives ─────────────────────────────────────────────────────────────

def _set_run_style(run, *, size=14, bold=False, color=FG, font=FONT_BODY, spacing=None):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if spacing is not None:
        run.font._rPr.set("spc", str(spacing))


def _add_textbox(slide, *, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, Emu(0))
    tf.vertical_anchor = anchor
    tf.paragraphs[0].text = ""
    return tf


def _add_paragraph(
    tf, text, *,
    size=14, bold=False, color=FG, font=FONT_BODY, spacing=None,
    space_before=0, space_after=4, align=PP_ALIGN.LEFT, first=False,
):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_run_style(run, size=size, bold=bold, color=color, font=font, spacing=spacing)
    return p


def _white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _accent_bar(slide, *, top, height_pt=4, width_in=0.6, left_in=0.7):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), top, Inches(width_in), Pt(height_pt),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _logo_corner(slide, *, size_in=0.35, right_pad=0.7, top_pad=0.35):
    """Small Mercor M logo in the top-right corner. Transparent PNG."""
    if not MERCOR_LOGO.exists():
        return
    slide_w_in = 13.333
    left = slide_w_in - right_pad - size_in
    slide.shapes.add_picture(
        str(MERCOR_LOGO),
        Inches(left), Inches(top_pad),
        Inches(size_in), Inches(size_in),
    )


def _footer(slide, *, page, total):
    tf = _add_textbox(slide, left=Inches(0.7), top=Inches(7.05), width=Inches(12), height=Inches(0.3))
    _add_paragraph(
        tf,
        f"creator experts   ·   internal proposal for aaron langerman   ·   {page} / {total}",
        size=9, color=FG_SUBTLE, font=FONT_BODY, first=True,
    )


def _micro_label(slide, *, text, top_in, left_in=0.7):
    tf = _add_textbox(slide, left=Inches(left_in), top=Inches(top_in), width=Inches(11), height=Inches(0.3))
    _add_paragraph(
        tf, text.upper(),
        size=10, bold=True, color=ACCENT, font=FONT_DISPLAY, spacing=120, first=True,
    )


def _hero(slide, *, text, top_in=1.5, size=44, color=FG, height=2.2,
          line_spacing=1.05, left_in=0.7, width_in=12):
    tf = _add_textbox(slide, left=Inches(left_in), top=Inches(top_in),
                      width=Inches(width_in), height=Inches(height))
    p = _add_paragraph(tf, text, size=size, bold=True, color=color, font=FONT_DISPLAY, first=True)
    p.line_spacing = line_spacing


def _rounded_card(slide, *, left_in, top_in, width_in, height_in,
                  fill_color=BG_ELEV, line_color=BORDER, line_pt=0.75,
                  corner_adj=0.06):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    box.adjustments[0] = corner_adj
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = line_color
    box.line.width = Pt(line_pt)
    return box


# ─── Slide builders ─────────────────────────────────────────────────────────

def slide_cover(prs, total):
    """Slide 1 — Cover. New thesis line per Aaron feedback."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_background(slide)
    _logo_corner(slide)

    _accent_bar(slide, top=Inches(2.0), width_in=0.8)
    _hero(slide, text="Creator Experts.", top_in=2.3, size=88, height=1.5, line_spacing=1.0)
    _hero(
        slide,
        text="A high-friction labor market. Where AI compounds.",
        top_in=4.0, size=24, color=FG_MUTED, height=0.7, line_spacing=1.05,
    )
    _hero(
        slide,
        text="Sourcing. Verification. Evaluation. Outcome prediction.",
        top_in=4.75, size=14, color=FG_SUBTLE, height=0.5, line_spacing=1.0,
    )

    # Byline block
    tf = _add_textbox(slide, left=Inches(0.7), top=Inches(5.9), width=Inches(12), height=Inches(1.0))
    _add_paragraph(
        tf, "For Aaron Langerman, Strategic Ops.",
        size=13, bold=True, color=FG, font=FONT_DISPLAY, first=True, space_after=14,
    )
    _add_paragraph(
        tf, "Logan Mann  +  Aarnav Nagabhirava",
        size=12, bold=True, color=FG, font=FONT_DISPLAY, space_after=2,
    )
    _add_paragraph(
        tf,
        "On Mercor's CUA-envs contract for Anthropic, under Drew Geoly.",
        size=11, color=FG_MUTED, font=FONT_BODY,
    )

    _footer(slide, page=1, total=total)


def slide_broken_workflow(prs, total):
    """Slide 2 — THE BROKEN WORKFLOW. NEW. 5-step horizontal flow + bottleneck stats."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_background(slide)
    _logo_corner(slide)
    _micro_label(slide, text="The broken workflow.", top_in=0.95)

    _hero(
        slide,
        text="Five steps. Most of them manual. None of them scale.",
        top_in=1.5, size=28, height=1.0, line_spacing=1.05,
    )

    # 5-step horizontal flow
    steps = [
        ("01", "Brand brief",   "marketing manager opens a doc"),
        ("02", "DM agents",     "searches IG/TikTok DMs"),
        ("03", "Manual scroll", "opens 40 profiles, evaluates by eye"),
        ("04", "Email volley",  "3–5 round-trips per creator"),
        ("05", "Sign or pass",  "guess at fit, sign contract, hope"),
    ]
    n = len(steps)
    margin = 0.7
    gap = 0.15
    avail = 13.333 - 2 * margin
    card_w = (avail - gap * (n - 1)) / n
    top = 3.0
    height = 1.85

    for i, (num, head, body) in enumerate(steps):
        left = margin + i * (card_w + gap)
        accent = i in (2,)  # highlight the manual-scroll bottleneck
        _rounded_card(
            slide, left_in=left, top_in=top, width_in=card_w, height_in=height,
            fill_color=ACCENT_SOFT if accent else BG_ELEV,
            line_color=ACCENT if accent else BORDER,
            line_pt=1.0 if accent else 0.75,
        )
        tf = _add_textbox(
            slide, left=Inches(left + 0.22), top=Inches(top + 0.22),
            width=Inches(card_w - 0.44), height=Inches(height - 0.4),
        )
        _add_paragraph(
            tf, num, size=11, bold=True,
            color=ACCENT if accent else FG_SUBTLE,
            font=FONT_DISPLAY, spacing=120, first=True, space_after=6,
        )
        _add_paragraph(
            tf, head, size=14, bold=True, color=FG, font=FONT_DISPLAY,
            space_after=4,
        )
        _add_paragraph(
            tf, body, size=10, color=FG_MUTED, font=FONT_BODY,
        )

    # Stats card row
    stats_top = 5.15
    stats_h = 1.45
    stats = [
        ("70%", "of brands say finding the right creators is their biggest bottleneck.",
         "Aspire, State of Influencer Marketing 2025."),
        ("39%", "of brands still rely on manual research.",
         "IMH/Sprout Q1 2025 Pulse."),
    ]
    sw = (avail - gap) / 2
    for i, (big, line, src) in enumerate(stats):
        left = margin + i * (sw + gap)
        _rounded_card(
            slide, left_in=left, top_in=stats_top, width_in=sw, height_in=stats_h,
            fill_color=BG_ELEV, line_color=BORDER,
        )
        tf = _add_textbox(
            slide, left=Inches(left + 0.3), top=Inches(stats_top + 0.22),
            width=Inches(sw - 0.6), height=Inches(stats_h - 0.4),
        )
        _add_paragraph(
            tf, big, size=30, bold=True, color=ACCENT,
            font=FONT_DISPLAY, first=True, space_after=4,
        )
        _add_paragraph(
            tf, line, size=12, bold=True, color=FG, font=FONT_BODY,
            space_after=4,
        )
        _add_paragraph(
            tf, src, size=9, color=FG_SUBTLE, font=FONT_BODY,
        )

    _footer(slide, page=2, total=total)


def slide_why_now(prs, total):
    """Slide 3 — WHY NOW. Three numbered drivers."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_background(slide)
    _logo_corner(slide)
    _micro_label(slide, text="Why now.", top_in=0.95)

    _hero(
        slide,
        text="Three forces just shifted under the creator market.",
        top_in=1.5, size=28, height=1.0,
    )

    drivers = [
        (
            "01",
            "Manual review doesn't scale when content explodes.",
            "YouTube enforcement, Jan 2026: 16 channels removed, 35M subs, "
            "4.7B views erased (Neal Mohan, 2026-01-12).",
        ),
        (
            "02",
            "Surface signals are getting noisier.",
            "Only 26% of consumers prefer AI-generated creator content today, "
            "down from 60% in 2023 (Billion Dollar Boy Muse Two, 2025-11-20).",
        ),
        (
            "03",
            "Ad costs are climbing.",
            "Meta Q3 2025: +10% YoY price-per-ad, "
            "\"increased advertiser demand\" (Susan Li, 2025-10-29).",
        ),
    ]
    top = 3.05
    row_h = 1.2
    for i, (num, head, body) in enumerate(drivers):
        y = top + i * row_h
        num_tf = _add_textbox(slide, left=Inches(0.7), top=Inches(y),
                              width=Inches(0.85), height=Inches(1))
        _add_paragraph(
            num_tf, num, size=24, bold=True, color=ACCENT,
            font=FONT_DISPLAY, first=True,
        )
        body_tf = _add_textbox(slide, left=Inches(1.7), top=Inches(y),
                               width=Inches(11), height=Inches(1.1))
        _add_paragraph(
            body_tf, head, size=18, bold=True, color=FG, font=FONT_DISPLAY,
            first=True, space_after=4,
        )
        _add_paragraph(
            body_tf, body, size=12, color=FG_MUTED, font=FONT_BODY,
        )

    _footer(slide, page=3, total=total)


def slide_signal(prs, total):
    """Slide 4 — SIGNAL. Comparison ledger + THE MOAT + 68% Lumanu anchor."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_background(slide)
    _logo_corner(slide)
    _micro_label(slide, text="Signal.", top_in=0.95)

    _hero(
        slide,
        text="Pay for proof, not for follower count.",
        top_in=1.5, size=28, height=1.0,
    )

    # Anchor stat strip — Lumanu 68% (replaces UNVERIFIED 78%)
    anchor_top = 2.55
    anchor_h = 0.85
    _rounded_card(
        slide, left_in=0.7, top_in=anchor_top, width_in=11.93, height_in=anchor_h,
        fill_color=ACCENT_SOFT, line_color=ACCENT, line_pt=1.0,
    )
    tf = _add_textbox(
        slide, left=Inches(0.95), top=Inches(anchor_top + 0.16),
        width=Inches(11.5), height=Inches(anchor_h - 0.3),
    )
    _add_paragraph(
        tf,
        "68% of brand-creator contracts include performance metrics — up from 42% in 2023.",
        size=15, bold=True, color=FG, font=FONT_DISPLAY, first=True, space_after=2,
    )
    _add_paragraph(
        tf, "Lumanu $1B+ payouts dataset, 2025.",
        size=10, color=FG_MUTED, font=FONT_BODY,
    )

    # Comparison ledger — header + 4 rows including THE MOAT
    ledger_top = 3.7
    col1_left = 0.7
    col1_w = 4.3
    col2_left = 5.2
    col2_w = 3.7
    col3_left = 9.05
    col3_w = 3.58
    row_h = 0.55
    rows = [
        ("Ranking input",   "Follower count",            "Outcome history × niche fit"),
        ("Trust check",     "Self-declared",             "Handle + post-fingerprint verified"),
        ("Pricing",         "Flat per post",             "Placement on close + perf kicker"),
        ("THE MOAT",        "—",                          "Predictor delta vs follower-baseline +18%, 90d"),
    ]

    # Header
    head_tf_a = _add_textbox(slide, left=Inches(col1_left), top=Inches(ledger_top),
                             width=Inches(col1_w), height=Inches(0.4))
    _add_paragraph(head_tf_a, "DIMENSION", size=9, bold=True, color=FG_SUBTLE,
                   font=FONT_DISPLAY, spacing=120, first=True)
    head_tf_b = _add_textbox(slide, left=Inches(col2_left), top=Inches(ledger_top),
                             width=Inches(col2_w), height=Inches(0.4))
    _add_paragraph(head_tf_b, "INDUSTRY DEFAULT", size=9, bold=True, color=FG_SUBTLE,
                   font=FONT_DISPLAY, spacing=120, first=True)
    head_tf_c = _add_textbox(slide, left=Inches(col3_left), top=Inches(ledger_top),
                             width=Inches(col3_w), height=Inches(0.4))
    _add_paragraph(head_tf_c, "MERCOR CREATOR DOMAIN", size=9, bold=True, color=ACCENT,
                   font=FONT_DISPLAY, spacing=120, first=True)

    # Rule line
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(col1_left), Inches(ledger_top + 0.32),
        Inches(col3_left + col3_w - col1_left), Pt(0.75),
    )
    rule.fill.solid(); rule.fill.fore_color.rgb = BORDER; rule.line.fill.background()

    base = ledger_top + 0.45
    for i, (dim, default, mercor) in enumerate(rows):
        y = base + i * row_h
        is_moat = dim == "THE MOAT"
        if is_moat:
            # subtle accent-soft band for the moat row
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(col1_left - 0.08), Inches(y - 0.05),
                Inches(col3_left + col3_w - col1_left + 0.16), Inches(row_h - 0.05),
            )
            band.fill.solid(); band.fill.fore_color.rgb = ACCENT_SOFT
            band.line.fill.background()
        tf_a = _add_textbox(slide, left=Inches(col1_left), top=Inches(y),
                            width=Inches(col1_w), height=Inches(row_h))
        _add_paragraph(
            tf_a, dim, size=11, bold=True,
            color=ACCENT if is_moat else FG, font=FONT_DISPLAY, first=True,
        )
        tf_b = _add_textbox(slide, left=Inches(col2_left), top=Inches(y),
                            width=Inches(col2_w), height=Inches(row_h))
        _add_paragraph(
            tf_b, default, size=11,
            color=FG_MUTED, font=FONT_BODY, first=True,
        )
        tf_c = _add_textbox(slide, left=Inches(col3_left), top=Inches(y),
                            width=Inches(col3_w), height=Inches(row_h))
        _add_paragraph(
            tf_c, mercor, size=11, bold=is_moat,
            color=FG, font=FONT_BODY, first=True,
        )

    _footer(slide, page=4, total=total)


def slide_verification(prs, total):
    """Slide 5 — VERIFICATION. NEW. 4-step ladder, steps 3+4 honest "manual review for v1"."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_background(slide)
    _logo_corner(slide)
    _micro_label(slide, text="Verification.", top_in=0.95)

    _hero(
        slide,
        text="Four checks to filter the fakes.",
        top_in=1.5, size=28, height=0.8,
    )

    # Left column — the ladder (4 steps)
    ladder_left = 0.7
    ladder_w = 7.5
    ladder_top = 2.55
    step_h = 0.95

    steps = [
        ("01", "Handle ownership",
         "TikTok/IG handle resolves to a live profile with matching display name.", False),
        ("02", "Fingerprint",
         "Recent post timestamps + caption style confirm consistent authorship.", False),
        ("03", "Niche claim",
         "Self-declared tags vs last 30 posts. Manual review for v1.", True),
        ("04", "Audience truth",
         "Sampled audience demographics check. Manual review for v1.", True),
    ]
    for i, (num, head, body, manual) in enumerate(steps):
        y = ladder_top + i * step_h
        # number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(ladder_left), Inches(y + 0.05),
            Inches(0.55), Inches(0.55),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_SOFT
        circle.line.color.rgb = ACCENT
        circle.line.width = Pt(0.75)
        ctf = circle.text_frame
        ctf.margin_left = Emu(0); ctf.margin_right = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.paragraphs[0].text = ""
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run(); crun.text = num
        _set_run_style(crun, size=12, bold=True, color=ACCENT, font=FONT_DISPLAY)

        body_tf = _add_textbox(
            slide, left=Inches(ladder_left + 0.75), top=Inches(y),
            width=Inches(ladder_w - 0.85), height=Inches(step_h),
        )
        _add_paragraph(
            body_tf, head, size=15, bold=True, color=FG, font=FONT_DISPLAY,
            first=True, space_after=2,
        )
        _add_paragraph(
            body_tf, body, size=11,
            color=FG_MUTED, font=FONT_BODY, space_after=2,
        )
        if manual:
            tag_tf = _add_textbox(
                slide, left=Inches(ladder_left + 0.75), top=Inches(y + 0.55),
                width=Inches(ladder_w - 0.85), height=Inches(0.3),
            )
            _add_paragraph(
                tag_tf, "MANUAL REVIEW FOR V1",
                size=8, bold=True, color=ACCENT, font=FONT_DISPLAY,
                spacing=140, first=True,
            )

    # Right column — three stat cards stacked
    stat_left = 8.55
    stat_w = 4.08
    stat_top = 2.55
    stat_h = 1.3
    stat_gap = 0.18
    stats = [
        ("37.2%", "of influencer followers are fake/inauthentic. ~$4.6B/yr brand waste.",
         "SociaVault 100K-account audit, 2025."),
        ("55%",   "of Instagram influencers have engaged in fraudulent activity.",
         "HypeAuditor State of Influencer Marketing 2024."),
        ("$53,088", "FTC max civil penalty per violation (Final Rule, eff. 2024-10-21).",
         "Federal Register, 2024."),
    ]
    for i, (big, line, src) in enumerate(stats):
        y = stat_top + i * (stat_h + stat_gap)
        _rounded_card(
            slide, left_in=stat_left, top_in=y, width_in=stat_w, height_in=stat_h,
            fill_color=BG_ELEV, line_color=BORDER,
        )
        tf = _add_textbox(
            slide, left=Inches(stat_left + 0.22), top=Inches(y + 0.16),
            width=Inches(stat_w - 0.44), height=Inches(stat_h - 0.3),
        )
        _add_paragraph(
            tf, big, size=22, bold=True, color=ACCENT,
            font=FONT_DISPLAY, first=True, space_after=2,
        )
        _add_paragraph(
            tf, line, size=10, color=FG, font=FONT_BODY, space_after=2,
        )
        _add_paragraph(
            tf, src, size=8, color=FG_SUBTLE, font=FONT_BODY,
        )

    _footer(slide, page=5, total=total)


def slide_day1_revenue(prs, total):
    """Slide 6 — DAY-1 REVENUE. Three lines on Logan × Celsius signed deal."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_background(slide)
    _logo_corner(slide)
    _micro_label(slide, text="Day-1 revenue.", top_in=0.95)

    _hero(
        slide,
        text="Three revenue lines. One signed deal.",
        top_in=1.5, size=28, height=0.9,
    )

    lines = [
        (
            "a", "Placement fee on close",
            "$170",
            "20% × $850 base contract.",
            "Headhunter contingent benchmark.",
        ),
        (
            "b", "Performance kicker (per-post)",
            "$0–$50",
            "5% of view-bonus pool.",
            "Lumanu pattern.",
        ),
        (
            "c", "RL data licensing",
            "Per task",
            "Outcome data + rubric scores feed Mercor's existing human-data market for RL environments.",
            "Illustrative — flagged in deck.",
        ),
    ]
    top = 2.7
    row_h = 1.05
    for i, (tag, head, dollars, body, src) in enumerate(lines):
        y = top + i * row_h
        # Tag pill
        pill = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.7), Inches(y + 0.1), Inches(0.45), Inches(0.45),
        )
        pill.fill.solid(); pill.fill.fore_color.rgb = ACCENT
        pill.line.fill.background()
        ctf = pill.text_frame
        ctf.margin_left = Emu(0); ctf.margin_right = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.paragraphs[0].text = ""
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run(); crun.text = tag
        _set_run_style(crun, size=14, bold=True, color=BG, font=FONT_DISPLAY)

        # Headline + body
        body_tf = _add_textbox(
            slide, left=Inches(1.35), top=Inches(y + 0.05),
            width=Inches(8.0), height=Inches(row_h),
        )
        _add_paragraph(
            body_tf, head, size=15, bold=True, color=FG, font=FONT_DISPLAY,
            first=True, space_after=2,
        )
        _add_paragraph(
            body_tf, body, size=11, color=FG_MUTED, font=FONT_BODY, space_after=2,
        )
        _add_paragraph(
            body_tf, src, size=9, color=FG_SUBTLE, font=FONT_BODY,
        )

        # Dollar column (right-aligned)
        money_tf = _add_textbox(
            slide, left=Inches(9.6), top=Inches(y + 0.15),
            width=Inches(3.05), height=Inches(0.7),
        )
        _add_paragraph(
            money_tf, dollars, size=22, bold=True, color=ACCENT,
            font=FONT_DISPLAY, align=PP_ALIGN.RIGHT, first=True,
        )

    # Bottom callout — verbatim DAY1_REVENUE_HEADLINE
    callout_top = 6.0
    callout_h = 0.9
    _rounded_card(
        slide, left_in=0.7, top_in=callout_top, width_in=11.93, height_in=callout_h,
        fill_color=ACCENT_SOFT, line_color=ACCENT, line_pt=1.0,
    )
    tf = _add_textbox(
        slide, left=Inches(0.95), top=Inches(callout_top + 0.18),
        width=Inches(11.5), height=Inches(callout_h - 0.3), anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_paragraph(
        tf, DAY1_REVENUE_HEADLINE,
        size=13, bold=True, color=FG, font=FONT_DISPLAY, first=True,
    )

    _footer(slide, page=6, total=total)


def slide_prove_by_august(prs, total):
    """Slide 7 — PROVE BY AUGUST. Four falsifiable proof points + verdict gate."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_background(slide)
    _logo_corner(slide)
    _micro_label(slide, text="Prove by august.", top_in=0.95)

    _hero(
        slide,
        text="Four proof points. Falsifiable. By August.",
        top_in=1.5, size=28, height=0.9,
    )

    # 2x2 grid of proof points
    proofs = [
        ("01", "Workflow pain real",
         "10 brand interviews. Quantify per-creator review time."),
        ("02", "Verification works",
         "1,000 creators screened. Catch rate vs ground truth."),
        ("03", "Quality scoring beats baseline",
         "Mercor model vs follower-count, on 5 brands."),
        ("04", "Outcome prediction holds",
         "Predicted vs actual on 10 finished campaigns."),
    ]
    grid_left = 0.7
    grid_top = 2.65
    cell_w = 5.95
    cell_h = 1.35
    cell_gap = 0.25
    for i, (num, head, body) in enumerate(proofs):
        col = i % 2
        row = i // 2
        x = grid_left + col * (cell_w + cell_gap)
        y = grid_top + row * (cell_h + cell_gap)
        _rounded_card(
            slide, left_in=x, top_in=y, width_in=cell_w, height_in=cell_h,
            fill_color=BG_ELEV, line_color=BORDER,
        )
        tf = _add_textbox(
            slide, left=Inches(x + 0.3), top=Inches(y + 0.2),
            width=Inches(cell_w - 0.6), height=Inches(cell_h - 0.35),
        )
        _add_paragraph(
            tf, num, size=10, bold=True, color=ACCENT,
            font=FONT_DISPLAY, spacing=120, first=True, space_after=4,
        )
        _add_paragraph(
            tf, head, size=15, bold=True, color=FG, font=FONT_DISPLAY,
            space_after=4,
        )
        _add_paragraph(
            tf, body, size=11, color=FG_MUTED, font=FONT_BODY,
        )

    # Verdict gate strip
    gate_top = 5.65
    gate_h = 0.85
    _rounded_card(
        slide, left_in=0.7, top_in=gate_top, width_in=11.93, height_in=gate_h,
        fill_color=ACCENT_SOFT, line_color=ACCENT, line_pt=1.0,
    )
    tf = _add_textbox(
        slide, left=Inches(0.95), top=Inches(gate_top + 0.16),
        width=Inches(11.5), height=Inches(gate_h - 0.3), anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_paragraph(
        tf, "4 of 4 → strategic adjacency. 2 of 4 → kill it.",
        size=15, bold=True, color=FG, font=FONT_DISPLAY, first=True, space_after=2,
    )
    _add_paragraph(
        tf, "Paid like a headhunter. Placement on signed deals.",
        size=11, color=FG_MUTED, font=FONT_BODY,
    )

    _footer(slide, page=7, total=total)


def slide_demo(prs, total):
    """Slide 8 — NOW → LIVE DEMO. Vercel URL handoff (same as v2 slide 7)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_background(slide)
    _logo_corner(slide)
    _micro_label(slide, text="Now → live demo", top_in=0.95)

    _hero(
        slide,
        text="The demo's the rest of the pitch.",
        top_in=1.55, size=42, height=2.3,
    )

    body_tf = _add_textbox(slide, left=Inches(0.7), top=Inches(4.4),
                           width=Inches(12), height=Inches(1.4))
    lines = [
        "Logan as a creator — real 22.7K TikTok account, applying inside Mercor's stepper.",
        "Aaron as Mercor admin — RAG cites real posts back by URL, outreach drafts in line.",
    ]
    first = True
    for body in lines:
        p = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = body
        _set_run_style(run, size=15, color=FG, font=FONT_BODY)
        first = False

    url_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(6.2), Inches(12), Inches(0.65),
    )
    url_box.adjustments[0] = 0.4
    url_box.fill.solid()
    url_box.fill.fore_color.rgb = ACCENT
    url_box.line.fill.background()
    url_tf = url_box.text_frame
    url_tf.margin_left = Emu(0)
    url_tf.margin_right = Emu(0)
    url_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    url_tf.paragraphs[0].text = ""
    p = url_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Live: musing-maxwell-84ed29.vercel.app"
    _set_run_style(run, size=15, bold=True, color=BG, font=FONT_DISPLAY)

    _footer(slide, page=8, total=total)


# ─── Build ──────────────────────────────────────────────────────────────────

def build_deck(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,
        slide_broken_workflow,
        slide_why_now,
        slide_signal,
        slide_verification,
        slide_day1_revenue,
        slide_prove_by_august,
        slide_demo,
    ]
    total = len(builders)
    for fn in builders:
        fn(prs, total)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"saved → {out_path}  ({total} slides)")


if __name__ == "__main__":
    here = Path(__file__).resolve().parent.parent
    out = here / "docs" / "pitch" / "Mercor-Creators-Domain.pptx"
    build_deck(out)
